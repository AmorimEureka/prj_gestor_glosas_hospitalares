from __future__ import annotations

import math
from pathlib import Path

import fitz
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
REPO = ROOT.parent
ASSETS = ROOT / "assets"
OUTPUT = ROOT / "Receita_Certa_Apresentacao_Alta_Gestao.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Aptos"

NAVY = "062751"
INK = "182230"
TEAL = "1F6F86"
TEAL_DARK = "164F63"
TEAL_LIGHT = "E7F4F8"
CYAN = "5FB6C7"
GOLD = "C88521"
ORANGE = "F3A65A"
GREEN = "227153"
GREEN_LIGHT = "E8F5EF"
RED = "B04444"
RED_LIGHT = "FCECEB"
BLUE_LIGHT = "E9EEF8"
BG = "F5F8FA"
WHITE = "FFFFFF"
GRAY = "667085"
GRAY_DARK = "475467"
LINE = "D0D5DD"
LINE_LIGHT = "E7EBEF"


def I(value: float):
    return Inches(value)


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor.from_string(value.upper())


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    font: str = FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.02,
    line_spacing: float = 1.0,
    italic: bool = False,
):
    box = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = I(margin)
    frame.margin_right = I(margin)
    frame.margin_top = I(margin)
    frame.margin_bottom = I(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    lines = text.split("\n")
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        paragraph.line_spacing = line_spacing
        paragraph.space_after = Pt(0)
        for run in paragraph.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = rgb(color)
    return box


def add_rich_text(
    slide,
    runs: list[tuple[str, dict]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = INK,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.02,
):
    box = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = I(margin)
    frame.margin_right = I(margin)
    frame.margin_top = I(margin)
    frame.margin_bottom = I(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    for value, style in runs:
        run = paragraph.add_run()
        run.text = value
        run.font.name = style.get("font", FONT)
        run.font.size = Pt(style.get("size", size))
        run.font.bold = style.get("bold", False)
        run.font.italic = style.get("italic", False)
        run.font.color.rgb = rgb(style.get("color", color))
    return box


def add_shape(
    slide,
    kind,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str | None = None,
    line_width: float = 1,
    radius: bool = False,
    rotation: float = 0,
):
    shape_kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else kind
    shape = slide.shapes.add_shape(shape_kind, I(x), I(y), I(w), I(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    shape.rotation = rotation
    return shape


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str | None = None,
    line_width: float = 1,
    radius: bool = False,
):
    return add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        x,
        y,
        w,
        h,
        fill=fill,
        line=line,
        line_width=line_width,
        radius=radius,
    )


def add_circle(
    slide,
    x: float,
    y: float,
    d: float,
    *,
    fill: str = WHITE,
    line: str | None = None,
    line_width: float = 1,
):
    return add_shape(
        slide,
        MSO_SHAPE.OVAL,
        x,
        y,
        d,
        d,
        fill=fill,
        line=line,
        line_width=line_width,
    )


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str = LINE,
    width: float = 1.5,
):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_arrow_right(
    slide,
    x1: float,
    y: float,
    x2: float,
    *,
    color: str = TEAL,
    width: float = 1.6,
):
    add_line(slide, x1, y, x2 - 0.10, y, color=color, width=width)
    tri = add_shape(
        slide,
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        x2 - 0.14,
        y - 0.07,
        0.14,
        0.14,
        fill=color,
        rotation=90,
    )
    return tri


def add_arrow_down(
    slide,
    x: float,
    y1: float,
    y2: float,
    *,
    color: str = TEAL,
    width: float = 1.6,
):
    add_line(slide, x, y1, x, y2 - 0.10, color=color, width=width)
    return add_shape(
        slide,
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        x - 0.07,
        y2 - 0.14,
        0.14,
        0.14,
        fill=color,
        rotation=180,
    )


def add_header(slide, title: str, kicker: str | None = None):
    if kicker:
        add_text(
            slide,
            kicker.upper(),
            0.72,
            0.40,
            6.8,
            0.24,
            size=8.5,
            color=TEAL,
            bold=True,
        )
    add_text(
        slide,
        title,
        0.72,
        0.70 if kicker else 0.48,
        11.9,
        0.62,
        size=27,
        color=NAVY,
        bold=True,
    )
    add_rect(slide, 0.72, 1.32, 0.70, 0.06, fill=TEAL)


def add_footer(slide, number: int, source: str | None = None):
    add_line(slide, 0.72, 7.13, 12.62, 7.13, color=LINE_LIGHT, width=0.7)
    if source:
        add_text(
            slide,
            source,
            0.74,
            7.19,
            10.8,
            0.16,
            size=6.5,
            color=GRAY,
        )
    add_text(
        slide,
        f"{number:02d}",
        12.15,
        7.17,
        0.42,
        0.18,
        size=7,
        color=GRAY,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def add_badge(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    *,
    fill: str = TEAL_LIGHT,
    color: str = TEAL_DARK,
):
    add_rect(slide, x, y, w, 0.31, fill=fill, radius=True)
    add_text(
        slide,
        text.upper(),
        x + 0.08,
        y + 0.075,
        w - 0.16,
        0.15,
        size=7.2,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str,
    *,
    accent: str = TEAL,
    fill: str = WHITE,
    title_size: float = 14,
    body_size: float = 10.8,
    label: str | None = None,
):
    add_rect(slide, x, y, w, h, fill=fill, line=LINE_LIGHT, radius=True)
    add_rect(slide, x, y, 0.08, h, fill=accent, radius=True)
    if label:
        add_text(
            slide,
            label.upper(),
            x + 0.24,
            y + 0.20,
            w - 0.40,
            0.17,
            size=7,
            color=accent,
            bold=True,
        )
        title_y = y + 0.46
    else:
        title_y = y + 0.25
    add_text(
        slide,
        title,
        x + 0.24,
        title_y,
        w - 0.43,
        0.45,
        size=title_size,
        color=NAVY,
        bold=True,
    )
    add_text(
        slide,
        body,
        x + 0.24,
        title_y + 0.55,
        w - 0.43,
        h - (title_y - y) - 0.70,
        size=body_size,
        color=GRAY_DARK,
        line_spacing=1.05,
    )


def add_metric_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    value: str,
    label: str,
    *,
    accent: str = TEAL,
    fill: str = WHITE,
    value_size: float = 22,
):
    add_rect(slide, x, y, w, h, fill=fill, line=LINE_LIGHT, radius=True)
    add_rect(slide, x, y, w, 0.055, fill=accent)
    add_text(
        slide,
        value,
        x + 0.20,
        y + 0.22,
        w - 0.40,
        0.42,
        size=value_size,
        color=NAVY,
        bold=True,
    )
    add_text(
        slide,
        label,
        x + 0.20,
        y + 0.73,
        w - 0.40,
        h - 0.86,
        size=9.2,
        color=GRAY,
    )


def add_numbered_step(
    slide,
    number: str,
    title: str,
    body: str,
    x: float,
    y: float,
    w: float,
    *,
    accent: str = TEAL,
):
    add_circle(slide, x, y, 0.42, fill=accent)
    add_text(
        slide,
        number,
        x,
        y + 0.105,
        0.42,
        0.16,
        size=8.5,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        title,
        x + 0.56,
        y - 0.01,
        w - 0.56,
        0.28,
        size=12.2,
        color=NAVY,
        bold=True,
    )
    add_text(
        slide,
        body,
        x + 0.56,
        y + 0.31,
        w - 0.56,
        0.50,
        size=9.2,
        color=GRAY,
    )


def prepare_assets():
    ASSETS.mkdir(parents=True, exist_ok=True)
    logo_svg = REPO / "frontend/static/img/prontocardio-slogan.svg"
    logo_png = ASSETS / "prontocardio.png"
    if not logo_png.exists():
        svg = logo_svg.read_bytes()
        doc = fitz.open(stream=svg, filetype="svg")
        pix = doc[0].get_pixmap(matrix=fitz.Matrix(2.5, 2.5), alpha=True)
        pix.save(logo_png)

    donut_png = ASSETS / "cobertura_zero_glosa.png"
    if not donut_png.exists():
        size = 900
        image = Image.new("RGBA", (size, size), (255, 255, 255, 0))
        draw = ImageDraw.Draw(image)
        box = (95, 95, size - 95, size - 95)
        width = 118
        draw.arc(box, -90, 360 - 90, fill="#" + ORANGE, width=width)
        covered_angle = 360 * 0.154
        draw.arc(
            box,
            -90,
            -90 + covered_angle,
            fill="#" + TEAL,
            width=width,
        )
        image.save(donut_png)
    return logo_png, donut_png


def new_slide(prs: Presentation, fill: str = BG):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = rgb(fill)
    return slide


def build_deck():
    logo_png, donut_png = prepare_assets()

    prs = Presentation()
    prs.slide_width = I(SLIDE_W)
    prs.slide_height = I(SLIDE_H)
    prs.core_properties.title = (
        "Receita Certa — implantação da plataforma integrada do ciclo da receita"
    )
    prs.core_properties.subject = (
        "Apresentação estratégica para a alta gestão do Hospital Prontocardio"
    )
    prs.core_properties.author = "Tecnologia da Informação — Hospital Prontocardio"
    prs.core_properties.comments = (
        "Conteúdo validado a partir dos repositórios api_prontocardio, "
        "prj_glosas e prj_web_nfs e do estudo interno da ZeroGlosa."
    )

    # 01 — CAPA
    slide = new_slide(prs, NAVY)
    add_rect(slide, 0, 0, 0.18, SLIDE_H, fill=CYAN)
    add_rect(slide, 8.92, 0, 4.41, SLIDE_H, fill=TEAL_DARK)
    add_circle(slide, 10.14, 0.18, 3.05, fill=TEAL)
    add_circle(slide, 9.45, 4.72, 2.35, fill=GOLD)
    add_circle(slide, 11.55, 5.60, 1.25, fill=CYAN)
    add_text(
        slide,
        "RECEITA CERTA",
        0.82,
        0.68,
        3.8,
        0.25,
        size=10,
        color=CYAN,
        bold=True,
    )
    add_text(
        slide,
        "Do indicador de glosa\nà orquestração do faturamento",
        0.82,
        1.25,
        7.35,
        1.45,
        size=29,
        color=WHITE,
        bold=True,
        line_spacing=0.92,
    )
    add_text(
        slide,
        "Implantação da plataforma integrada do ciclo da receita",
        0.84,
        2.98,
        6.55,
        0.45,
        size=15,
        color="C8D7E7",
    )
    add_text(
        slide,
        "Visão estratégica • funcionalidades • problemas resolvidos • "
        "comparativo ZeroGlosa",
        0.84,
        3.52,
        6.85,
        0.42,
        size=10.5,
        color="AFC6D4",
    )

    flow_y = 5.18
    stages = [
        ("PRODUÇÃO", 0.88, 1.08),
        ("FISCAL", 2.33, 0.84),
        ("CAIXA", 3.56, 0.84),
        ("GLOSAS", 4.79, 0.94),
        ("RECUPERAÇÃO", 6.14, 1.42),
    ]
    for index, (label, x, w) in enumerate(stages):
        add_badge(
            slide,
            label,
            x,
            flow_y,
            w,
            fill="173F64",
            color=WHITE,
        )
        if index < len(stages) - 1:
            add_arrow_right(
                slide,
                x + w + 0.07,
                flow_y + 0.15,
                stages[index + 1][1] - 0.08,
                color=CYAN,
                width=1.4,
            )
    add_text(
        slide,
        "Alta Gestão | Julho de 2026",
        0.84,
        6.72,
        3.4,
        0.20,
        size=8,
        color="AFC6D4",
        bold=True,
    )
    slide.shapes.add_picture(str(logo_png), I(10.08), I(1.92), width=I(2.15))

    # 02 — TESE EXECUTIVA
    slide = new_slide(prs)
    add_header(slide, "A tese executiva em uma frase", "Síntese")
    add_rect(slide, 0.72, 1.72, 11.90, 1.38, fill=NAVY, radius=True)
    add_text(
        slide,
        "“A glosa era o sintoma. O problema real era a fragmentação "
        "do ciclo da receita.”",
        1.08,
        2.05,
        11.15,
        0.65,
        size=23,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    outcomes = [
        (
            "01",
            "Visão fim a fim",
            "Produção hospitalar, documento fiscal, recebimento, glosa e recuperação no mesmo fluxo.",
            TEAL,
        ),
        (
            "02",
            "Automação crítica",
            "A emissão de NFS-e deixa de depender de e-mail, planilha e digitação repetitiva no portal.",
            GOLD,
        ),
        (
            "03",
            "Gestão por evidência",
            "Indicadores mostram valor, prazo, convênio, motivo, recurso, perda e recuperação.",
            GREEN,
        ),
        (
            "04",
            "Governança operacional",
            "Permissões, trilhas de auditoria, estados de processo e tratamento de exceções.",
            NAVY,
        ),
    ]
    for i, (num, title, body, accent) in enumerate(outcomes):
        x = 0.72 + i * 3.02
        add_rect(slide, x, 3.56, 2.83, 2.12, fill=WHITE, line=LINE_LIGHT, radius=True)
        add_circle(slide, x + 0.22, 3.80, 0.46, fill=accent)
        add_text(
            slide,
            num,
            x + 0.22,
            3.93,
            0.46,
            0.14,
            size=7.5,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            title,
            x + 0.82,
            3.78,
            1.78,
            0.34,
            size=13,
            color=NAVY,
            bold=True,
        )
        add_text(
            slide,
            body,
            x + 0.23,
            4.43,
            2.34,
            0.88,
            size=10,
            color=GRAY_DARK,
            line_spacing=1.05,
        )
    add_rect(slide, 1.62, 6.08, 10.08, 0.58, fill=TEAL_LIGHT, radius=True)
    add_text(
        slide,
        "TI como alavanca de resultado operacional — não apenas como suporte.",
        1.92,
        6.27,
        9.48,
        0.22,
        size=12.5,
        color=TEAL_DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 2)

    # 03 — INFLEXÃO ESTRATÉGICA
    slide = new_slide(prs)
    add_header(slide, "Da demanda inicial à visão sistêmica", "Inflexão estratégica")
    add_line(slide, 1.20, 3.12, 8.14, 3.12, color=LINE, width=2.2)
    timeline = [
        (
            1.28,
            TEAL,
            "Demanda inicial",
            "Gerar indicadores de glosas hospitalares.",
        ),
        (
            3.95,
            GOLD,
            "Sinal da Direção",
            "TI como força estratégica para otimizar a máquina administrativa de faturamento.",
        ),
        (
            6.68,
            GREEN,
            "Decisão de produto",
            "Tratar o ciclo da receita como sistema integrado — do fato gerador ao caixa.",
        ),
    ]
    for index, (x, accent, title, body) in enumerate(timeline, 1):
        add_circle(slide, x, 2.82, 0.60, fill=accent)
        add_text(
            slide,
            f"{index}",
            x,
            3.00,
            0.60,
            0.18,
            size=10,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            title,
            x - 0.44,
            3.70,
            1.55,
            0.42,
            size=11.5,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            body,
            x - 0.64,
            4.18,
            2.0,
            1.18,
            size=9.7,
            color=GRAY_DARK,
            align=PP_ALIGN.CENTER,
        )
    add_rect(slide, 9.18, 1.75, 3.44, 4.88, fill=NAVY, radius=True)
    add_badge(
        slide,
        "Leitura estratégica",
        9.58,
        2.12,
        1.92,
        fill="173F64",
        color=CYAN,
    )
    add_text(
        slide,
        "Não bastava medir\nonde a receita se perde.",
        9.58,
        2.72,
        2.64,
        0.86,
        size=16,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        "Era necessário redesenhar como a receita é produzida, documentada, recebida, "
        "contestada e recuperada.",
        9.58,
        3.82,
        2.62,
        1.28,
        size=11,
        color="C8D7E7",
        line_spacing=1.08,
    )
    add_line(slide, 9.58, 5.40, 11.82, 5.40, color="456A85", width=1)
    add_text(
        slide,
        "Resultado: Receita Certa",
        9.58,
        5.70,
        2.62,
        0.30,
        size=12,
        color=CYAN,
        bold=True,
    )
    add_text(
        slide,
        "Uma plataforma, não uma tela isolada.",
        9.58,
        6.08,
        2.62,
        0.24,
        size=9.5,
        color="C8D7E7",
    )
    add_footer(
        slide,
        3,
        "Fonte: contexto estratégico informado e evidências funcionais dos três repositórios.",
    )

    # 04 — PROBLEMAS ANTES
    slide = new_slide(prs)
    add_header(
        slide,
        "Antes: o ciclo existia, mas não operava como sistema",
        "Diagnóstico",
    )
    process = [
        ("PRODUÇÃO", "MV / ERP", NAVY),
        ("PEDIDO DE NOTA", "E-mail ou planilha", RED),
        ("EMISSÃO FISCAL", "Portal manual", RED),
        ("RECEBIMENTO", "Controles separados", GOLD),
        ("GLOSAS", "Histórico fragmentado", RED),
    ]
    y = 2.15
    for i, (title, subtitle, accent) in enumerate(process):
        x = 0.74 + i * 2.48
        add_rect(slide, x, y, 2.03, 1.23, fill=WHITE, line=LINE, radius=True)
        add_rect(slide, x, y, 2.03, 0.07, fill=accent)
        add_text(
            slide,
            title,
            x + 0.16,
            y + 0.25,
            1.72,
            0.20,
            size=8.2,
            color=accent,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            subtitle,
            x + 0.12,
            y + 0.65,
            1.80,
            0.48,
            size=10.3,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        if i < len(process) - 1:
            add_line(
                slide,
                x + 2.08,
                y + 0.62,
                x + 2.39,
                y + 0.62,
                color=RED,
                width=2,
            )
            add_text(
                slide,
                "×",
                x + 2.17,
                y + 0.40,
                0.18,
                0.24,
                size=15,
                color=RED,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
    add_text(
        slide,
        "Fragmentação operacional",
        0.78,
        3.72,
        2.42,
        0.28,
        size=12,
        color=RED,
        bold=True,
    )
    problems = [
        ("Retrabalho", "O mesmo dado era digitado e conferido em múltiplos pontos."),
        ("Tempo de ciclo", "Pedidos, validação, emissão e retorno dependiam de passagem manual."),
        ("Risco fiscal", "Falhas de cadastro, duplicidade e baixa rastreabilidade da emissão."),
        ("Receita invisível", "Produção, NFS-e, recebimento e glosa não fechavam a mesma conta."),
        ("Gestão reativa", "Sem visão integrada de prazo, causa, responsável e recuperação."),
    ]
    for i, (title, body) in enumerate(problems):
        x = 0.74 + i * 2.48
        add_circle(slide, x + 0.03, 4.36, 0.34, fill=RED_LIGHT, line=RED)
        add_text(
            slide,
            "!",
            x + 0.03,
            4.45,
            0.34,
            0.14,
            size=8,
            color=RED,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            title,
            x + 0.48,
            4.32,
            1.58,
            0.24,
            size=11,
            color=NAVY,
            bold=True,
        )
        add_text(
            slide,
            body,
            x + 0.48,
            4.72,
            1.64,
            1.14,
            size=8.7,
            color=GRAY_DARK,
        )
    add_rect(slide, 2.04, 6.17, 9.25, 0.53, fill=RED_LIGHT, radius=True)
    add_text(
        slide,
        "O impacto não era apenas produtividade: era atraso de caixa, perda de evidência e risco de receita.",
        2.30,
        6.34,
        8.72,
        0.25,
        size=10,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 4, "Fonte: processo anterior descrito pelo solicitante.")

    # 05 — ARQUITETURA EXECUTIVA
    slide = new_slide(prs)
    add_header(slide, "Receita Certa: uma camada de integração e controle", "Solução")

    # Usuários
    user_cards = [
        ("RECEPÇÕES", "Clínicas e emergência"),
        ("FINANCEIRO", "Valida, emite e concilia"),
        ("GLOSAS", "Trata, recorre e acompanha"),
        ("GESTÃO", "Decide por indicadores"),
    ]
    for i, (title, subtitle) in enumerate(user_cards):
        x = 0.78 + i * 2.38
        add_rect(slide, x, 1.78, 2.08, 0.78, fill=WHITE, line=LINE_LIGHT, radius=True)
        add_text(
            slide,
            title,
            x + 0.12,
            1.98,
            1.84,
            0.16,
            size=8.5,
            color=TEAL,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            subtitle,
            x + 0.12,
            2.23,
            1.84,
            0.16,
            size=7.5,
            color=GRAY,
            align=PP_ALIGN.CENTER,
        )
        add_arrow_down(slide, x + 1.04, 2.58, 2.88, color=TEAL, width=1.2)

    add_rect(slide, 0.78, 2.93, 9.24, 1.35, fill=NAVY, radius=True)
    add_text(
        slide,
        "EXPERIÊNCIA E ORQUESTRAÇÃO",
        1.10,
        3.18,
        2.62,
        0.20,
        size=8,
        color=CYAN,
        bold=True,
    )
    add_text(
        slide,
        "prj_glosas",
        1.10,
        3.54,
        2.62,
        0.34,
        size=17,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        "Interface web • fluxos por perfil • indicadores • acompanhamento",
        3.42,
        3.34,
        5.97,
        0.34,
        size=11,
        color="C8D7E7",
        align=PP_ALIGN.RIGHT,
    )

    add_rect(slide, 10.38, 1.78, 2.22, 2.50, fill=TEAL_LIGHT, line="CFE4E8", radius=True)
    add_text(
        slide,
        "GESTÃO DO CICLO",
        10.66,
        2.05,
        1.66,
        0.18,
        size=8,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Produção\n→ Fiscal\n→ Caixa\n→ Glosa\n→ Recuperação",
        10.66,
        2.48,
        1.66,
        1.38,
        size=13,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        line_spacing=0.95,
    )

    add_arrow_down(slide, 5.40, 4.30, 4.62, color=TEAL, width=1.5)
    add_rect(slide, 0.78, 4.67, 9.24, 0.90, fill=WHITE, line=LINE, radius=True)
    add_text(
        slide,
        "INTEGRAÇÃO E REGRAS",
        1.08,
        4.93,
        2.32,
        0.18,
        size=8,
        color=GOLD,
        bold=True,
    )
    add_text(
        slide,
        "api_prontocardio",
        3.12,
        4.87,
        2.52,
        0.28,
        size=15,
        color=NAVY,
        bold=True,
    )
    add_text(
        slide,
        "Dados • autenticação • conciliação • auditoria • workflow",
        5.30,
        4.91,
        4.25,
        0.24,
        size=9.5,
        color=GRAY_DARK,
        align=PP_ALIGN.RIGHT,
    )

    systems = [
        ("MV / ORACLE", "Produção hospitalar", NAVY),
        ("POSTGRESQL", "Estado e auditoria", TEAL),
        ("AIRFLOW", "prj_web_nfs", GOLD),
        ("ISS FORTALEZA", "Emissão / XML / PDF", GREEN),
    ]
    for i, (title, subtitle, accent) in enumerate(systems):
        x = 0.78 + i * 2.95
        add_arrow_down(slide, x + 1.25, 5.59, 5.86, color=accent, width=1.2)
        add_rect(slide, x, 5.90, 2.50, 0.83, fill=WHITE, line=LINE_LIGHT, radius=True)
        add_text(
            slide,
            title,
            x + 0.15,
            6.12,
            2.20,
            0.16,
            size=8.4,
            color=accent,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            subtitle,
            x + 0.15,
            6.39,
            2.20,
            0.16,
            size=7.4,
            color=GRAY,
            align=PP_ALIGN.CENTER,
        )
    add_footer(
        slide,
        5,
        "Fonte: arquitetura observada em prj_glosas, api_prontocardio e prj_web_nfs.",
    )

    # 06 — CAPACIDADES
    slide = new_slide(prs)
    add_header(slide, "Quatro capacidades que mudam a operação", "Funcionalidades")
    cards = [
        (
            0.72,
            1.72,
            "01",
            "Faturamento e conciliação",
            "• Produção hospitalar × NFS-e\n• Remessas × notas fiscais\n• Previsão e registro de recebimento\n• Extrato bancário, impostos e retenções\n• Consulta e edição auditável",
            TEAL,
        ),
        (
            6.83,
            1.72,
            "02",
            "Gestão de glosas",
            "• Triagem e follow-up por remessa\n• Recurso, acato e recebimento\n• Motivos padronizados pela TISS\n• Prazo por convênio e aging\n• Visão Kanban e tabela",
            GOLD,
        ),
        (
            0.72,
            4.28,
            "03",
            "Inteligência gerencial",
            "• Funil financeiro por convênio\n• Pareto dos motivos de glosa\n• Eficiência de recuperação\n• Evolução mensal e SLA\n• Filtros por período, prestador e atendimento",
            GREEN,
        ),
        (
            6.83,
            4.28,
            "04",
            "Automação fiscal",
            "• Pedido digital pelas recepções\n• Validação pelo contas a receber\n• Emissão individual ou em lote\n• Número, protocolo e PDF na aplicação\n• Erro visível, corrigível e reprocessável",
            NAVY,
        ),
    ]
    for x, y, num, title, body, accent in cards:
        add_rect(slide, x, y, 5.78, 2.15, fill=WHITE, line=LINE_LIGHT, radius=True)
        add_circle(slide, x + 0.25, y + 0.28, 0.60, fill=accent)
        add_text(
            slide,
            num,
            x + 0.25,
            y + 0.46,
            0.60,
            0.17,
            size=9,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            title,
            x + 1.05,
            y + 0.28,
            4.35,
            0.38,
            size=16,
            color=NAVY,
            bold=True,
        )
        add_text(
            slide,
            body,
            x + 1.05,
            y + 0.82,
            4.35,
            1.08,
            size=10.1,
            color=GRAY_DARK,
            line_spacing=1.04,
        )
    add_footer(
        slide,
        6,
        "Fonte: telas, rotas, modelos e regras de negócio observados nos três repositórios.",
    )

    # 07 — NFS-e
    slide = new_slide(prs)
    add_header(slide, "NFS-e: resposta rápida a um ponto crítico do ERP", "Automação")
    add_rect(slide, 0.72, 1.70, 3.28, 4.98, fill=RED_LIGHT, line="F1D0CD", radius=True)
    add_badge(
        slide,
        "Situação de ruptura",
        1.08,
        2.05,
        1.96,
        fill="F7DCD9",
        color=RED,
    )
    add_text(
        slide,
        "A rotina de emissão pelo ERP deixou de atender a operação.",
        1.08,
        2.62,
        2.57,
        0.92,
        size=17,
        color=NAVY,
        bold=True,
    )
    add_text(
        slide,
        "O contingenciamento passou a depender de pedido por e-mail ou planilha "
        "e emissão manual no portal da Prefeitura.",
        1.08,
        3.83,
        2.57,
        1.04,
        size=11,
        color=GRAY_DARK,
    )
    add_rect(slide, 1.08, 5.18, 2.57, 0.94, fill=WHITE, radius=True)
    add_text(
        slide,
        "Risco imediato",
        1.28,
        5.40,
        2.17,
        0.18,
        size=8,
        color=RED,
        bold=True,
    )
    add_text(
        slide,
        "Atraso, retrabalho e dependência operacional.",
        1.28,
        5.67,
        2.17,
        0.34,
        size=9.5,
        color=NAVY,
        bold=True,
    )

    add_text(
        slide,
        "Fluxo implantado",
        4.48,
        1.78,
        2.58,
        0.28,
        size=14,
        color=TEAL,
        bold=True,
    )
    steps = [
        ("1", "Solicitação digital", "Clínicas e emergência registram o pedido."),
        ("2", "Validação financeira", "Dados, valor e emissor são conferidos."),
        ("3", "Seleção do lote", "Itens aprovados seguem para processamento."),
        ("4", "Robô fiscal", "Portal da Prefeitura é preenchido e validado."),
        ("5", "Retorno automático", "Número, protocolo e PDF voltam à aplicação."),
        ("6", "Gestão por exceção", "Erros ficam visíveis para correção e nova tentativa."),
    ]
    for i, (num, title, body) in enumerate(steps):
        col = 0 if i < 3 else 1
        row = i if i < 3 else i - 3
        x = 4.48 + col * 4.05
        y = 2.30 + row * 1.17
        accent = TEAL if col == 0 else GREEN
        add_numbered_step(slide, num, title, body, x, y, 3.64, accent=accent)
    add_rect(slide, 4.47, 5.93, 8.13, 0.74, fill=NAVY, radius=True)
    add_text(
        slide,
        "Pessoas validam e decidem. O robô executa a repetição e devolve a evidência.",
        4.78,
        6.16,
        7.52,
        0.26,
        size=12,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(
        slide,
        7,
        "Fonte: contexto informado; workflow, estados, anti-duplicidade e integração ISS observados no código.",
    )

    # 08 — INDICADORES
    slide = new_slide(prs)
    add_header(slide, "Indicadores que conectam operação e resultado", "Inteligência gerencial")
    kpis = [
        ("TOTAL GLOSADO", TEAL),
        ("NÃO TRATADAS", RED),
        ("RECURSADAS", GOLD),
        ("ACATADAS", GRAY),
        ("RECUPERADAS", GREEN),
        ("RECURSOS EM ABERTO", NAVY),
    ]
    for i, (label, accent) in enumerate(kpis):
        x = 0.72 + i * 2.0
        add_rect(slide, x, 1.72, 1.82, 0.85, fill=WHITE, line=LINE_LIGHT, radius=True)
        add_rect(slide, x, 1.72, 1.82, 0.055, fill=accent)
        add_text(
            slide,
            "R$ —",
            x + 0.16,
            1.99,
            1.50,
            0.20,
            size=12,
            color=NAVY,
            bold=True,
        )
        add_text(
            slide,
            label,
            x + 0.16,
            2.32,
            1.50,
            0.14,
            size=6.6,
            color=accent,
            bold=True,
        )

    # Funil conceitual
    add_rect(slide, 0.72, 2.88, 7.30, 3.28, fill=WHITE, line=LINE_LIGHT, radius=True)
    add_text(
        slide,
        "Funil financeiro por convênio",
        1.02,
        3.17,
        3.32,
        0.26,
        size=13.5,
        color=NAVY,
        bold=True,
    )
    funnel = [
        ("Faturado", 5.86, TEAL),
        ("Glosado", 4.76, GOLD),
        ("Recursado", 3.76, NAVY),
        ("Recuperado", 2.66, GREEN),
        ("Perda", 1.72, RED),
    ]
    for i, (label, width, color) in enumerate(funnel):
        x = 1.06 + (5.86 - width) / 2
        y = 3.70 + i * 0.43
        add_rect(slide, x, y, width, 0.30, fill=color, radius=True)
        add_text(
            slide,
            label,
            x,
            y + 0.075,
            width,
            0.12,
            size=6.8,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    add_text(
        slide,
        "Conversão entre etapas evidencia onde a receita se perde e onde a atuação retorna caixa.",
        1.08,
        5.82,
        6.58,
        0.22,
        size=8.5,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )

    # Quadrante conceitual
    add_rect(slide, 8.32, 2.88, 4.30, 3.28, fill=WHITE, line=LINE_LIGHT, radius=True)
    add_text(
        slide,
        "Eficiência por motivo de glosa",
        8.62,
        3.17,
        3.70,
        0.26,
        size=13.5,
        color=NAVY,
        bold=True,
    )
    qx, qy, qw, qh = 8.82, 3.70, 3.30, 1.82
    add_rect(slide, qx, qy, qw / 2, qh / 2, fill=GREEN_LIGHT)
    add_rect(slide, qx + qw / 2, qy, qw / 2, qh / 2, fill="DDF2E8")
    add_rect(slide, qx, qy + qh / 2, qw / 2, qh / 2, fill="F5F6F7")
    add_rect(slide, qx + qw / 2, qy + qh / 2, qw / 2, qh / 2, fill=RED_LIGHT)
    add_line(slide, qx + qw / 2, qy, qx + qw / 2, qy + qh, color=LINE)
    add_line(slide, qx, qy + qh / 2, qx + qw, qy + qh / 2, color=LINE)
    quadrant_labels = [
        ("Baixa prioridade", qx + 0.12, qy + 0.12, GREEN),
        ("Excelente", qx + 1.78, qy + 0.12, GREEN),
        ("Pouco relevante", qx + 0.12, qy + 1.03, GRAY),
        ("Prioridade máxima", qx + 1.78, qy + 1.03, RED),
    ]
    for label, x, y, color in quadrant_labels:
        add_text(slide, label, x, y, 1.38, 0.17, size=6.8, color=color, bold=True)
    bubbles = [
        (9.42, 4.27, 0.22, GOLD),
        (10.90, 4.04, 0.36, GREEN),
        (11.17, 4.92, 0.48, RED),
        (9.64, 5.08, 0.17, GRAY),
    ]
    for x, y, d, color in bubbles:
        add_circle(slide, x, y, d, fill=color, line=WHITE, line_width=1)
    add_text(
        slide,
        "Filtros: período • tipo • convênio • prestador • atendimento • motivo",
        8.62,
        5.82,
        3.70,
        0.34,
        size=7.8,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    add_rect(slide, 2.48, 6.40, 8.38, 0.42, fill=TEAL_LIGHT, radius=True)
    add_text(
        slide,
        "A gestão deixa de perguntar “quanto perdemos?” e passa a decidir “onde agir primeiro?”.",
        2.70,
        6.53,
        7.94,
        0.17,
        size=10.5,
        color=TEAL_DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(
        slide,
        8,
        "Visual conceitual baseado nos indicadores implementados; valores dependem do período e filtros selecionados.",
    )

    # 09 — ZERO GLOSA
    slide = new_slide(prs)
    add_header(slide, "ZeroGlosa: o que os dados realmente demonstram", "Comparativo econômico")
    slide.shapes.add_picture(str(donut_png), I(0.86), I(1.76), width=I(3.62))
    add_text(
        slide,
        "15,4%",
        1.67,
        2.82,
        1.98,
        0.46,
        size=26,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "da receita avaliada\ncoberta",
        1.67,
        3.34,
        1.98,
        0.42,
        size=9,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "R$ 13,44 mi cobertos",
        1.16,
        5.26,
        1.66,
        0.24,
        size=8.6,
        color=TEAL,
        bold=True,
    )
    add_text(
        slide,
        "R$ 74,00 mi fora do escopo",
        2.66,
        5.26,
        2.00,
        0.24,
        size=8.6,
        color=GOLD,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )

    metrics = [
        (4.80, 1.82, "R$ 104,4 mil", "custo bruto acumulado em 17 meses", RED),
        (8.65, 1.82, "12,9%", "das contas em remessa cobertas", GOLD),
        (4.80, 3.22, "R$ 6,14 mil", "custo médio mensal histórico", TEAL),
        (8.65, 3.22, "R$ 5,83", "custo por conta coberta", TEAL),
        (4.80, 4.62, "0,78%", "break-even sobre a receita coberta", GREEN),
        (8.65, 4.62, "R$ 74,0 mi", "faturamento fora do escopo", NAVY),
    ]
    for x, y, value, label, accent in metrics:
        add_metric_card(slide, x, y, 3.48, 1.10, value, label, accent=accent, value_size=20)
    add_rect(slide, 4.80, 6.04, 7.33, 0.60, fill=RED_LIGHT, radius=True)
    add_text(
        slide,
        "Ponto central: não havia histórico integrado para atribuir ao fornecedor "
        "recursos e valores efetivamente recuperados.",
        5.08,
        6.14,
        6.76,
        0.40,
        size=10,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(
        slide,
        9,
        "Fonte: apresentação_custo_beneficio_zero_glosa.pdf — dados internos atualizados em 23/06/2026. Cálculos derivados: custo/17 meses e custo/17.893 contas.",
    )

    # 10 — POSICIONAMENTO COMPARATIVO
    slide = new_slide(prs)
    add_header(slide, "A comparação correta é de modelo operacional", "Receita Certa × ZeroGlosa")
    x0, y0 = 0.72, 1.72
    widths = [2.50, 4.12, 5.28]
    headers = [
        ("CRITÉRIO", NAVY),
        ("ZERO GLOSA", GOLD),
        ("RECEITA CERTA", TEAL),
    ]
    current_x = x0
    for (label, color), width in zip(headers, widths):
        add_rect(slide, current_x, y0, width, 0.56, fill=color)
        add_text(
            slide,
            label,
            current_x + 0.15,
            y0 + 0.19,
            width - 0.30,
            0.16,
            size=8.5,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        current_x += width
    rows = [
        (
            "Escopo",
            "Atuação sobre uma parcela das glosas.",
            "Produção → fiscal → caixa → glosa → recuperação.",
        ),
        (
            "Cobertura",
            "15,4% da receita; 12,9% das contas.",
            "Multiconvênio; cobertura depende de dados e rollout.",
        ),
        (
            "Evidência de retorno",
            "Recursos e recuperações não integrados ao histórico.",
            "Recurso, acato, recebimento e sucesso no modelo de dados.",
        ),
        (
            "Automação operacional",
            "NFS-e não avaliada no estudo.",
            "Pedido → validação → emissão → PDF → nova tentativa.",
        ),
        (
            "Governança",
            "Dependência do escopo e de evidência externa.",
            "Dados institucionais, permissões, auditoria e estados.",
        ),
        (
            "Leitura econômica",
            "R$ 104,4 mil / 17 meses; efetividade a demonstrar.",
            "TCO interno a medir para fechar ROI e contrato.",
        ),
    ]
    row_h = 0.68
    for i, (criterion, vendor, internal) in enumerate(rows):
        y = y0 + 0.56 + i * row_h
        fill = WHITE if i % 2 == 0 else "F9FAFB"
        current_x = x0
        for width in widths:
            add_rect(slide, current_x, y, width, row_h, fill=fill, line=LINE_LIGHT)
            current_x += width
        add_text(
            slide,
            criterion,
            x0 + 0.20,
            y + 0.22,
            widths[0] - 0.40,
            0.24,
            size=9.7,
            color=NAVY,
            bold=True,
        )
        add_text(
            slide,
            vendor,
            x0 + widths[0] + 0.22,
            y + 0.14,
            widths[1] - 0.44,
            0.38,
            size=9.2,
            color=GRAY_DARK,
        )
        add_text(
            slide,
            internal,
            x0 + widths[0] + widths[1] + 0.22,
            y + 0.14,
            widths[2] - 0.44,
            0.38,
            size=9.2,
            color=GRAY_DARK,
        )
    add_rect(slide, 1.55, 6.58, 10.23, 0.35, fill=TEAL_LIGHT, radius=True)
    add_text(
        slide,
        "A plataforma interna amplia o campo de gestão; a decisão financeira exige medir TCO, adoção e recuperação atribuível.",
        1.75,
        6.65,
        9.83,
        0.20,
        size=8.6,
        color=TEAL_DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(
        slide,
        10,
        "Comparação baseada no estudo interno da ZeroGlosa e nas capacidades observadas nos repositórios. Não presume cobertura integral nem economia já realizada.",
    )

    # 11 — CENÁRIOS
    slide = new_slide(prs)
    add_header(slide, "A escala muda a lógica econômica", "Potencial de valor")
    add_rect(slide, 0.72, 1.72, 4.06, 4.98, fill=NAVY, radius=True)
    add_badge(
        slide,
        "Base fora do fornecedor",
        1.10,
        2.10,
        1.96,
        fill="173F64",
        color=CYAN,
    )
    add_text(
        slide,
        "R$ 74,0 mi",
        1.10,
        2.72,
        3.30,
        0.62,
        size=30,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        "de faturamento fora do escopo\nno período avaliado",
        1.10,
        3.46,
        2.92,
        0.68,
        size=12,
        color="C8D7E7",
    )
    add_line(slide, 1.10, 4.48, 4.32, 4.48, color="456A85", width=1)
    add_text(
        slide,
        "0,14%",
        1.10,
        4.85,
        1.50,
        0.46,
        size=24,
        color=CYAN,
        bold=True,
    )
    add_text(
        slide,
        "de recuperação ou perda evitada nessa base equivale aproximadamente "
        "ao custo bruto de R$ 104,4 mil do fornecedor.",
        2.25,
        4.83,
        2.14,
        1.10,
        size=9.8,
        color=WHITE,
    )

    add_text(
        slide,
        "Cenários sobre a base histórica fora do escopo",
        5.22,
        1.92,
        6.72,
        0.30,
        size=12.5,
        color=NAVY,
        bold=True,
    )
    scenarios = [
        ("0,1%", 74_000, TEAL, "R$ 74 mil"),
        ("0,5%", 370_000, GOLD, "R$ 370 mil"),
        ("1,0%", 740_000, GREEN, "R$ 740 mil"),
    ]
    max_value = max(item[1] for item in scenarios)
    for i, (rate, value, color, label) in enumerate(scenarios):
        y = 2.62 + i * 1.08
        add_text(
            slide,
            rate,
            5.22,
            y + 0.10,
            0.62,
            0.24,
            size=12,
            color=color,
            bold=True,
            align=PP_ALIGN.RIGHT,
        )
        add_rect(slide, 6.08, y, 5.20, 0.48, fill="E9EEF1", radius=True)
        bar_width = 5.20 * value / max_value
        add_rect(slide, 6.08, y, bar_width, 0.48, fill=color, radius=True)
        add_text(
            slide,
            label,
            11.47,
            y + 0.10,
            0.92,
            0.22,
            size=10,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.RIGHT,
        )
    add_rect(slide, 5.22, 5.85, 7.18, 0.72, fill=RED_LIGHT, radius=True)
    add_text(
        slide,
        "Cenários de sensibilidade — não representam benefício realizado.",
        5.52,
        6.06,
        6.58,
        0.22,
        size=10.5,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Além do caixa recuperado: tempo de ciclo, horas operacionais, erros fiscais e continuidade do processo.",
        5.22,
        4.98,
        7.18,
        0.42,
        size=10,
        color=GRAY_DARK,
        italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(
        slide,
        11,
        "Fonte: base de R$ 74,00 mi do estudo ZeroGlosa. Cálculos: 0,1%, 0,5%, 1,0% e R$ 104,4 mil ÷ R$ 74,0 mi = 0,141%.",
    )

    # 12 — IMPLANTAÇÃO E GOVERNANÇA
    slide = new_slide(prs)
    add_header(slide, "Implantação com base técnica e agenda de valor", "Governança")
    add_text(
        slide,
        "Evidências do produto construído",
        0.72,
        1.70,
        5.64,
        0.30,
        size=14,
        color=TEAL,
        bold=True,
    )
    evidence = [
        ("3", "repositórios integrados", "Interface, API/dados e automação fiscal."),
        ("292", "cenários de teste no código", "101 no frontend, 155 na API e 36 na automação."),
        ("RBAC", "acesso por perfil e tela", "Separação de funções entre operação, gestão e TI."),
        ("AUDIT", "trilhas e estados rastreáveis", "Alterações, eventos, lotes, erros, PDFs e protocolos."),
    ]
    for i, (value, title, body) in enumerate(evidence):
        col, row = i % 2, i // 2
        x = 0.72 + col * 2.92
        y = 2.18 + row * 1.65
        is_code = len(value) > 3
        value_size = 15 if is_code else 18
        title_offset = 1.18 if is_code else 0.92
        add_rect(slide, x, y, 2.66, 1.40, fill=WHITE, line=LINE_LIGHT, radius=True)
        add_text(
            slide,
            value,
            x + 0.20,
            y + 0.22,
            0.94 if is_code else 0.78,
            0.34,
            size=value_size,
            color=NAVY,
            bold=True,
        )
        add_text(
            slide,
            title,
            x + title_offset,
            y + 0.24,
            2.44 - title_offset,
            0.34,
            size=9.4,
            color=TEAL,
            bold=True,
        )
        add_text(
            slide,
            body,
            x + 0.20,
            y + 0.78,
            2.22,
            0.42,
            size=8.1,
            color=GRAY_DARK,
        )
    add_text(
        slide,
        "Próximos 90 dias: transformar capacidade em resultado medido",
        6.78,
        1.70,
        5.84,
        0.30,
        size=14,
        color=GOLD,
        bold=True,
    )
    next_steps = [
        ("01", "Baseline", "Congelar linha de base: tempo de emissão, glosa, recurso, recebimento e custo."),
        ("02", "Adoção", "Formalizar responsáveis, SLAs e uso do fluxo único por área."),
        ("03", "Resultado", "Publicar painel mensal de recuperação, aging, perdas e produtividade."),
        ("04", "Decisão", "Comparar TCO interno e recuperação atribuível antes da decisão contratual."),
    ]
    for i, (num, title, body) in enumerate(next_steps):
        y = 2.18 + i * 0.95
        add_numbered_step(slide, num, title, body, 6.78, y, 5.45, accent=GOLD)
    add_rect(slide, 0.72, 5.92, 11.90, 0.75, fill=TEAL_LIGHT, radius=True)
    add_text(
        slide,
        "Métrica executiva recomendada: dinheiro recuperado + tempo até o caixa + custo operacional do ciclo.",
        1.08,
        6.17,
        11.18,
        0.24,
        size=12,
        color=TEAL_DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(
        slide,
        12,
        "Contagem de testes baseada nas funções de teste catalogadas em 29/07/2026; não equivale a relatório de execução nesta apresentação.",
    )

    # 13 — DECISÃO EXECUTIVA
    slide = new_slide(prs, NAVY)
    add_rect(slide, 0, 0, 0.18, SLIDE_H, fill=CYAN)
    add_text(
        slide,
        "DECISÃO EXECUTIVA",
        0.82,
        0.55,
        3.2,
        0.22,
        size=9,
        color=CYAN,
        bold=True,
    )
    add_text(
        slide,
        "O que precisamos da Alta Gestão",
        0.82,
        0.98,
        8.6,
        0.62,
        size=27,
        color=WHITE,
        bold=True,
    )
    decisions = [
        (
            "01",
            "Patrocinar o Receita Certa",
            "Reconhecer a solução como plataforma corporativa do ciclo da receita.",
        ),
        (
            "02",
            "Definir donos e SLAs",
            "Recepção, faturamento, financeiro, glosas e TI com responsabilidade explícita.",
        ),
        (
            "03",
            "Medir valor mensalmente",
            "Tempo, custo, glosa, recuperação, perdas, exceções e adoção do fluxo.",
        ),
        (
            "04",
            "Reavaliar o fornecedor por evidência",
            "Contrato comparado a recuperação atribuível, cobertura e TCO institucional.",
        ),
    ]
    for i, (num, title, body) in enumerate(decisions):
        col, row = i % 2, i // 2
        x = 0.82 + col * 5.72
        y = 2.02 + row * 1.55
        add_rect(slide, x, y, 5.35, 1.22, fill="173F64", line="315C78", radius=True)
        add_circle(slide, x + 0.22, y + 0.24, 0.48, fill=TEAL)
        add_text(
            slide,
            num,
            x + 0.22,
            y + 0.38,
            0.48,
            0.15,
            size=8,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            title,
            x + 0.86,
            y + 0.23,
            4.13,
            0.28,
            size=13,
            color=WHITE,
            bold=True,
        )
        add_text(
            slide,
            body,
            x + 0.86,
            y + 0.63,
            4.13,
            0.34,
            size=9,
            color="C8D7E7",
        )
    add_line(slide, 0.82, 5.43, 12.48, 5.43, color="315C78", width=1)
    add_text(
        slide,
        "A visão estratégica não foi automatizar uma tarefa.",
        0.82,
        5.80,
        7.85,
        0.40,
        size=18,
        color=CYAN,
        bold=True,
    )
    add_text(
        slide,
        "Foi criar disciplina institucional sobre a receita.",
        0.82,
        6.26,
        7.85,
        0.40,
        size=18,
        color=WHITE,
        bold=True,
    )
    slide.shapes.add_picture(str(logo_png), I(10.42), I(5.63), width=I(1.58))
    add_text(
        slide,
        "Receita Certa",
        10.01,
        6.83,
        2.40,
        0.20,
        size=8,
        color="AFC6D4",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    output = build_deck()
    print(output)
